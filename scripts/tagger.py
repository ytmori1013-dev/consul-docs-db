"""
PPT/PDFファイルのタグ付けモジュール（pdfplumberベース視覚密度検出版）

- PDF: pdfplumber の page.images / page.rects / page.lines を使って
  横長スライド判定とビジュアル密度判定を行う
- PPT/PPTX: python-pptx で図形数・テキスト量を集計
- 質評価はしない。客観的事実（密度・形式）のみをタグに反映する
- 抽出失敗時もタイトルで必ずテーマタグを付与して返す
- html/unknown エントリはタイトル・年度からテーマ・年度タグのみ付与（DLせず）
"""
import logging
import os
import re
import tempfile
from typing import Optional

import requests

try:
    from scripts.firms import detect_firm
except ImportError:
    from firms import detect_firm

logger = logging.getLogger(__name__)

# ===== 定数 =====
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
MAX_BATCH = 2000  # tag_entries() が処理する最大件数

HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/120.0.0.0 Safari/537.36"
    )
}

# ===== 構造タグ定義 =====
STRUCTURE_TAG_RULES = {
    "Issue Tree": ["論点", "課題ツリー", "イシュー", "Issue", "論点整理", "課題構造"],
    "MECE": ["MECE", "漏れなく", "ダブりなく", "網羅的", "切り口"],
    "So What": ["示唆", "インプリケーション", "So What", "したがって", "含意"],
    "ピラミッド構造": ["結論から", "キーメッセージ", "主張", "論拠"],
    "仮説思考": ["仮説", "Hypothesis", "仮説検証", "検証"],
    "ファクトベース": ["データ分析", "定量", "統計", "エビデンス", "実績値"],
}

# ===== テーマタグ定義 =====
THEME_TAG_RULES = {
    "DX・デジタル": ["DX", "デジタル", "Digital", "AI", "IoT", "クラウド", "データ活用"],
    "人的資本": ["人材", "HR", "人的資本", "採用", "育成", "スキル", "リスキリング", "雇用"],
    "GX・脱炭素": ["GX", "カーボン", "脱炭素", "ESG", "再生可能エネルギー", "水素", "EV"],
    "スタートアップ": ["スタートアップ", "ベンチャー", "イノベーション", "新規事業", "VC"],
    "社会保障": ["医療", "介護", "年金", "社会保障", "少子化", "高齢化"],
    "産業政策": ["産業政策", "製造業", "サプライチェーン", "半導体", "経済安全保障"],
    "地域・まちづくり": ["地方創生", "まちづくり", "自治体", "地域活性", "観光", "農業"],
    "海外・グローバル": ["海外展開", "グローバル", "輸出", "国際競争", "ASEAN"],
    "宇宙・防衛": ["宇宙", "衛星", "SDA", "SSA", "コンステ", "防衛", "安全保障", "ミサイル", "JAXA"],
}

# 注目スライドのタイトルキーワード
HIGHLIGHT_TITLE_KW = ["全体像", "サマリ", "まとめ", "ロードマップ", "フレームワーク"]


# ===== ファイルダウンロード =====

def download_file(url: str) -> Optional[str]:
    """
    URLからファイルをダウンロードして一時ファイルパスを返す。
    .pdf/.ppt/.pptx のみ対象。失敗・サイズ超過時は None。
    タイムアウト30秒・最大3回リトライ。
    """
    url_lower = url.lower().split("?")[0]
    if not any(url_lower.endswith(ext) for ext in (".pdf", ".ppt", ".pptx")):
        return None

    import time as _time
    for attempt in range(3):
        try:
            resp = requests.get(url, headers=HEADERS, timeout=30, stream=True)
            resp.raise_for_status()

            content_length = int(resp.headers.get("Content-Length", 0))
            if content_length > MAX_FILE_SIZE:
                logger.warning(f"ファイルサイズ超過({content_length}バイト)スキップ: {url}")
                return None

            suffix = (
                ".pptx" if url_lower.endswith(".pptx") else
                ".ppt"  if url_lower.endswith(".ppt")  else ".pdf"
            )
            total = 0
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as f:
                for chunk in resp.iter_content(chunk_size=65536):
                    f.write(chunk)
                    total += len(chunk)
                    if total > MAX_FILE_SIZE:
                        logger.warning(f"ダウンロード中にサイズ超過: {url}")
                        f.close()
                        os.unlink(f.name)
                        return None
                return f.name
        except Exception as e:
            if attempt < 2:
                _time.sleep(2 ** attempt)
            else:
                logger.warning(f"ダウンロード失敗: {url}: {e}")
    return None


# ===== PDF 抽出（pdfplumberベース） =====

def extract_from_pdf(file_path: str) -> dict:
    """
    pdfplumber で PDF を解析し、ビジュアル密度・テキスト密度を集計する。

    返却値:
      - slide_count: 総ページ数
      - landscape_ratio: 横長ページの割合
      - slide_type: "slide" (landscape_ratio>=0.7) / "document"
      - avg_visuals: (img+大きな矩形) の1ページ平均
      - avg_chars: 文字数の1ページ平均
      - full_text: 全ページ結合テキスト
      - per_slide: [{page, visual_count, title}, ...]
    """
    result = {
        "slide_count": 0,
        "landscape_ratio": 0.0,
        "slide_type": "document",
        "avg_visuals": 0.0,
        "avg_chars": 0.0,
        "full_text": "",
        "per_slide": [],
    }
    try:
        import pdfplumber

        landscape_pages = 0
        total_visuals = 0
        total_chars = 0
        per_slide = []
        all_texts = []

        with pdfplumber.open(file_path) as pdf:
            pages = pdf.pages
            n = len(pages)
            if n == 0:
                return result
            result["slide_count"] = n

            for idx, page in enumerate(pages, start=1):
                # 横長判定
                w = page.width or 0
                h = page.height or 0
                if w > h:
                    landscape_pages += 1

                # ビジュアル要素数の集計
                # page.rects は表のセル枠線も含むため面積フィルタをかける
                # (幅×高さ > 2000pt² ≒ 45×45pt 以上のみ「図形」として計上)
                img_count = len(page.images) if hasattr(page, "images") else 0
                large_rects = sum(
                    1 for r in (page.rects if hasattr(page, "rects") else [])
                    if (r.get("width", 0) * r.get("height", 0)) > 2000
                )
                visual_count = img_count + large_rects

                # テキスト抽出
                text = page.extract_text() or ""
                char_count = len(text)

                total_visuals += visual_count
                total_chars   += char_count
                all_texts.append(text)

                # タイトル：先頭行
                title = ""
                stripped = text.strip()
                if stripped:
                    title = stripped.split("\n")[0][:60]

                per_slide.append({
                    "page": idx,
                    "visual_count": visual_count,
                    "title": title,
                })

        # 横長スライド判定
        landscape_ratio = landscape_pages / n
        result["landscape_ratio"] = round(landscape_ratio, 4)
        result["slide_type"] = "slide" if landscape_ratio >= 0.7 else "document"

        # 密度指標
        result["avg_visuals"] = round(total_visuals / n, 2)
        result["avg_chars"]   = round(total_chars   / n, 2)
        result["full_text"]   = "\n".join(all_texts)
        result["per_slide"]   = per_slide

    except Exception as e:
        logger.warning(f"pdfplumber 抽出失敗: {e}")

    return result


# ===== PPT/PPTX 抽出 =====

def extract_from_pptx(file_path: str) -> dict:
    """
    python-pptx で PPT/PPTX を解析し、ビジュアル密度・テキスト密度を集計する。

    PPT/PPTX はスライド形式確定なので:
      - slide_type = "slide"
      - landscape_ratio = 1.0

    返却値:
      - slide_count, landscape_ratio, slide_type
      - avg_visuals: (PICTURE+AUTO_SHAPE+GROUP+CHART+TABLE 図形総数) / スライド数
      - avg_chars: テキスト総文字数 / スライド数
      - full_text, per_slide
    """
    result = {
        "slide_count": 0,
        "landscape_ratio": 1.0,
        "slide_type": "slide",
        "avg_visuals": 0.0,
        "avg_chars": 0.0,
        "full_text": "",
        "per_slide": [],
    }
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # カウント対象の図形タイプ
        VISUAL_TYPES = {
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.GROUP,
            MSO_SHAPE_TYPE.CHART,
            MSO_SHAPE_TYPE.TABLE,
        }

        prs = Presentation(file_path)
        n = len(prs.slides)
        if n == 0:
            return result

        total_visuals = 0
        total_chars   = 0
        all_texts     = []
        per_slide     = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_visuals = 0
            slide_texts   = []
            slide_title   = ""

            for shape in slide.shapes:
                # ビジュアル図形カウント
                try:
                    if shape.shape_type in VISUAL_TYPES:
                        slide_visuals += 1
                except Exception:
                    pass

                # テーブルのテキスト
                if getattr(shape, "has_table", False):
                    try:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                t = (cell.text or "").strip()
                                if t:
                                    slide_texts.append(t)
                    except Exception:
                        pass

                # テキストフレームのテキスト
                if getattr(shape, "has_text_frame", False):
                    try:
                        for para in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in para.runs).strip()
                            if line:
                                slide_texts.append(line)
                                if not slide_title:
                                    slide_title = line
                    except Exception:
                        pass

            slide_text = "\n".join(slide_texts)
            char_count = len(slide_text)

            total_visuals += slide_visuals
            total_chars   += char_count
            all_texts.append(slide_text)

            per_slide.append({
                "page": idx,
                "visual_count": slide_visuals,
                "title": slide_title[:60],
            })

        result["slide_count"] = n
        result["avg_visuals"] = round(total_visuals / n, 2)
        result["avg_chars"]   = round(total_chars   / n, 2)
        result["full_text"]   = "\n".join(all_texts)
        result["per_slide"]   = per_slide

    except Exception as e:
        logger.warning(f"PPTX抽出失敗: {e}")

    return result


# ===== タグ判定 =====

def detect_design_tags(avg_visuals: Optional[float], avg_chars: Optional[float]) -> list:
    """
    avg_visuals / avg_chars からデザインタグを判定する（客観的事実のみ）。
    avg_visuals / avg_chars が None の場合はタグを付与しない。
    """
    if avg_visuals is None or avg_chars is None:
        return []

    tags = []

    # 主タグ（3択・必ず1つ）
    # avg_visuals: 画像 + 大きな矩形の1スライド平均（表枠線は除外済み）
    # avg_chars: 文字数の1スライド平均
    if avg_visuals >= 3 and avg_chars < 150:
        tags.append("図解中心")
    elif avg_chars >= 250:
        tags.append("テキスト重視")
    else:
        tags.append("図表バランス型")

    # ビジュアル密度（必須）
    if avg_visuals >= 4:
        tags.append("ビジュアル密度：高")
    elif avg_visuals >= 2:
        tags.append("ビジュアル密度：中")
    else:
        tags.append("ビジュアル密度：低")

    # 情報密度（必須）
    if avg_chars >= 300:
        tags.append("情報密度：高")
    elif avg_chars >= 100:
        tags.append("情報密度：中")
    else:
        tags.append("情報密度：低")

    return tags


def score_structure_tags(text: str) -> list:
    """
    構造タグをキーワード出現回数でスコア化する。
    出現回数 >= 3 で ◎、1-2 で ○、0は付与しない。
    返却値: [{"name": "Issue Tree", "level": "◎"}, ...]
    """
    result = []
    for tag_name, keywords in STRUCTURE_TAG_RULES.items():
        score = sum(text.count(kw) for kw in keywords)
        if score >= 3:
            result.append({"name": tag_name, "level": "◎"})
        elif score >= 1:
            result.append({"name": tag_name, "level": "○"})
    return result


def detect_theme_tags(text: str) -> list:
    """テキストからテーマタグを検出する"""
    return [name for name, kws in THEME_TAG_RULES.items() if any(kw in text for kw in kws)]


def detect_year_tag(title: str, fiscal_year: str) -> list:
    """年度タグを生成する"""
    if fiscal_year:
        return [fiscal_year]
    m = re.search(r"(令和\d+|平成\d+)", title)
    if m:
        return [m.group(1)]
    return []


def detect_highlight_slides(per_slide: list) -> list:
    """
    注目スライドを最大3ページ抽出する（ページ番号の配列）。
    タイトルキーワード含むページ優先、次いでビジュアル数の多いページ上位3件。
    """
    if not per_slide:
        return []

    scored = []
    for s in per_slide:
        visual_count = s.get("visual_count", 0)
        title = s.get("title", "")

        # タイトルキーワードマッチを優先スコアとして加算
        score = visual_count
        if any(kw in title for kw in HIGHLIGHT_TITLE_KW):
            score += 1000  # キーワードマッチを最優先

        scored.append((score, s["page"]))

    scored.sort(reverse=True)
    return [p for _, p in scored[:3]]


# ===== メイン処理 =====

def tag_entry(entry: dict) -> dict:
    """
    1件のエントリーにタグを付与して返す。

    - PDF: pdfplumber で横長判定・ビジュアル密度を集計
    - PPT/PPTX: python-pptx で図形数・テキスト量を集計
    - html/unknown: タイトル・NDL メタデータでテーマ・年度タグを付与（DLせず）
    - 抽出失敗時もタイトルで必ずテーマタグを付与して返す
    """
    url         = entry.get("url", "")
    title       = entry.get("title", "")
    file_type   = entry.get("file_type", "")
    fiscal_year = entry.get("fiscal_year", "")
    firm_name   = entry.get("firm_name", "不明")

    # 返却値の初期化
    slide_count     = entry.get("slide_count", 0)
    landscape_ratio = None
    slide_type      = entry.get("slide_type", "")
    avg_visuals     = None
    avg_chars       = None
    per_slide       = []
    extraction_failed = entry.get("extraction_failed", False)
    full_text       = title  # 最低限タイトルは含める
    file_path       = None

    # html/unknown エントリ: NDL メタデータをテキスト判定に活用
    if file_type in ("html", "unknown"):
        ndl_meta = " ".join(filter(None, [
            entry.get("ndl_subject", ""),
            entry.get("ndl_description", ""),
            entry.get("ndl_responsibility", ""),
        ]))
        full_text = title + " " + ndl_meta
        if firm_name == "不明" and ndl_meta:
            firm_name = detect_firm(ndl_meta)

    try:
        if file_type in ("pdf", "ppt", "pptx"):
            file_path = download_file(url)
            if file_path:
                # ファイル種別ごとの抽出
                if file_type in ("ppt", "pptx"):
                    ext = extract_from_pptx(file_path)
                else:
                    ext = extract_from_pdf(file_path)

                extracted_text  = ext.get("full_text", "")
                slide_count     = ext.get("slide_count", 0) or slide_count
                landscape_ratio = ext.get("landscape_ratio", None)
                slide_type      = ext.get("slide_type", "") or slide_type
                avg_visuals     = ext.get("avg_visuals", None)
                avg_chars       = ext.get("avg_chars", None)
                per_slide       = ext.get("per_slide", [])

                logger.info(
                    f"抽出: {len(extracted_text)}文字 slides={slide_count} "
                    f"avg_visuals={avg_visuals} avg_chars={avg_chars} ({title[:30]})"
                )

                # テキストが極端に少ない場合は抽出失敗とみなす（ファームはタイトルのみ利用）
                if extracted_text.strip():
                    full_text = title + "\n" + extracted_text
                    if firm_name == "不明":
                        firm_name = detect_firm(extracted_text[:3000])
                else:
                    extraction_failed = True
            else:
                extraction_failed = True

    except Exception as e:
        logger.warning(f"テキスト抽出スキップ ({url}): {e}")
        extraction_failed = True
    finally:
        if file_path:
            try:
                os.unlink(file_path)
            except Exception:
                pass

    # ファーム検出（抽出できなかった場合はタイトルから試みる）
    if firm_name == "不明":
        firm_name = detect_firm(full_text)

    # PPT/PPTX の landscape_ratio はファイル解析で確定済み（None → デフォルト 1.0）
    if file_type in ("ppt", "pptx") and landscape_ratio is None:
        landscape_ratio = 1.0
        slide_type = "slide"

    # デザインタグ：avg_visuals/avg_chars が None（html等）の場合は付与しない
    design_tags = detect_design_tags(avg_visuals, avg_chars)

    tagged = {
        **entry,
        "firm_name": firm_name,
        "slide_type": slide_type,
        "landscape_ratio": landscape_ratio,
        "page_count": slide_count,
        "avg_visuals": avg_visuals,
        "avg_chars": avg_chars,
        "extraction_failed": extraction_failed,
        "highlight_slides": detect_highlight_slides(per_slide),
        "tags": {
            "structure": score_structure_tags(full_text),
            "design": design_tags,
            "theme": detect_theme_tags(full_text),
            "year": detect_year_tag(title, fiscal_year),
        },
    }
    return tagged


def tag_entries(entries: list) -> list:
    """
    エントリーリストにタグを付与する。
    先頭 MAX_BATCH 件のみ処理する。
    """
    targets = entries[:MAX_BATCH]
    total   = len(targets)
    results = []

    for i, entry in enumerate(targets):
        try:
            logger.info(f"タグ付け中 ({i + 1}/{total}): {entry.get('title', '')[:50]}")
            results.append(tag_entry(entry))
        except Exception as e:
            logger.error(f"タグ付けエラー（空タグで保存）: {e}")
            results.append({
                **entry,
                "firm_name": entry.get("firm_name", "不明"),
                "slide_type": entry.get("slide_type", ""),
                "landscape_ratio": None,
                "page_count": entry.get("slide_count", 0),
                "avg_visuals": None,
                "avg_chars": None,
                "extraction_failed": True,
                "highlight_slides": [],
                "tags": {"structure": [], "design": [], "theme": [], "year": []},
            })

    return results


if __name__ == "__main__":
    import json
    import sys
    logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
    entries = json.load(sys.stdin)
    print(json.dumps(tag_entries(entries), ensure_ascii=False, indent=2))
